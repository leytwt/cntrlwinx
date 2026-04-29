from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os


def create_pptx(slides_data):
    # Если будет template.pptx — можно подключить его
    # prs = Presentation("template.pptx")

    prs = Presentation()

    # -----------------------------
    # Титульный слайд
    # -----------------------------
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])

    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "AI Generated Presentation"
    subtitle.text = "Автоматически создано с помощью LLM"

    # -----------------------------
    # Основные слайды
    # -----------------------------
    for slide_data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data["title"]


        text_frame = content.text_frame
        text_frame.clear()

        points = slide_data["content"].split(". ")

        for point in points:
            if point.strip():
                p = text_frame.add_paragraph()
                p.text = point.strip()
                p.font.size = Pt(18)



    final_slide = prs.slides.add_slide(prs.slide_layouts[1])

    final_slide.shapes.title.text = "Спасибо!"

    final_content = final_slide.placeholders[1]
    final_content.text = "Готовы ответить на вопросы"


    if not os.path.exists("generated"):
        os.makedirs("generated")

    filename = f"generated/presentation_{datetime.now().timestamp()}.pptx"

    prs.save(filename)

    return filename