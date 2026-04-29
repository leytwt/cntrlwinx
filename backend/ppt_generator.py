# backend/ppt_generator.py

from pptx import Presentation
from image_service import ImageService
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os


# =====================================
# PREMIUM DESIGN SYSTEM
# =====================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# default = dark premium
DARK_BG = RGBColor(12, 18, 32)
CARD_BG = RGBColor(22, 28, 45)
WHITE = RGBColor(255, 255, 255)
SOFT_TEXT = RGBColor(203, 213, 225)
LIGHT_LINE = RGBColor(59, 130, 246)

# default theme
DEFAULT_ACCENT = RGBColor(59, 130, 246)


# =====================================
# COLOR RESOLVER
# =====================================

def resolve_theme_from_prompt(prompt: str):
    """
    User controls colors:
    pink / blue / black / gold etc.
    If nothing specified → premium dark default
    """

    text = prompt.lower()

    if "розов" in text or "pink" in text:
        return {
            "bg": RGBColor(35, 18, 35),
            "accent": RGBColor(236, 72, 153),
            "title": WHITE,
            "text": RGBColor(244, 214, 230)
        }

    if "син" in text or "blue" in text:
        return {
            "bg": RGBColor(12, 18, 32),
            "accent": RGBColor(59, 130, 246),
            "title": WHITE,
            "text": SOFT_TEXT
        }

    if "фиолет" in text or "purple" in text:
        return {
            "bg": RGBColor(24, 18, 42),
            "accent": RGBColor(168, 85, 247),
            "title": WHITE,
            "text": RGBColor(220, 210, 255)
        }

    if "золот" in text or "gold" in text or "luxury" in text:
        return {
            "bg": RGBColor(18, 18, 18),
            "accent": RGBColor(212, 175, 55),
            "title": WHITE,
            "text": RGBColor(235, 226, 200)
        }

    # default = premium dark
    return {
        "bg": DARK_BG,
        "accent": DEFAULT_ACCENT,
        "title": WHITE,
        "text": SOFT_TEXT
    }


# =====================================
# BASE BLOCKS
# =====================================

def add_background(slide, prs, bg_color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height
    )

    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()

    # send to back
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_accent_line(slide, accent_color):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(0.55),
        Inches(1.8),
        Inches(0.08)
    )

    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.background()


def split_content(content):
    if ";" in content:
        return [x.strip() for x in content.split(";") if x.strip()]

    if "\n" in content:
        return [x.strip() for x in content.split("\n") if x.strip()]

    return [content]


# =====================================
# SAFE ZONES (NO OVERLAPS)
# =====================================

def add_title(
    slide,
    text,
    title_color,
    image_on_left=False
):
    """
    Hard safe-zone:
    title NEVER touches image
    """

    left = Inches(6.7) if image_on_left else Inches(0.8)

    # if title is too long → make image smaller later
    width = Inches(5.0)

    box = slide.shapes.add_textbox(
        left,
        Inches(0.8),
        width,
        Inches(1.4)
    )

    tf = box.text_frame
    p = tf.paragraphs[0]

    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = title_color


def add_content(
    slide,
    content,
    text_color,
    image_on_left=False
):
    """
    HARD SAFE TEXT ZONE
    Текст никогда не уходит под картинку
    + нормальный перенос строк
    """

    left = Inches(6.7) if image_on_left else Inches(0.8)

    # делаем уже зону текста
    width = Inches(4.4)

    box = slide.shapes.add_textbox(
        left,
        Inches(2.0),
        width,
        Inches(4.2)
    )

    tf = box.text_frame

    # ВАЖНО:
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    points = split_content(content)

    for i, point in enumerate(points[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        # без ручного bullet символа
        # PowerPoint сам переносит лучше
        p.text = point

        p.font.size = Pt(16)
        p.font.color.rgb = text_color
        p.space_after = Pt(14)
        p.level = 0


def add_image_card(
    slide,
    image_path,
    image_on_left=False,
    smaller=False
):
    """
    Premium rounded image card
    визуально выглядит как rounded image
    """

    left = Inches(0.8) if image_on_left else Inches(6.4)

    card_width = Inches(4.8 if smaller else 5.3)
    card_height = Inches(3.8 if smaller else 4.2)

    top = Inches(1.2)

    # внешняя rounded card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        card_width,
        card_height
    )

    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(28, 35, 55)
    card.line.fill.background()

    # внутренние отступы
    padding = Inches(0.18)

    image_left = left + padding
    image_top = top + padding
    image_width = card_width - (padding * 2)
    image_height = card_height - (padding * 2)

    # картинка идеально внутри
    slide.shapes.add_picture(
        image_path,
        image_left,
        image_top,
        width=image_width,
        height=image_height
    )

    # мягкая overlay-рамка
    border = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left,
        top,
        card_width,
        card_height
    )

    border.fill.background()
    border.line.color.rgb = RGBColor(45, 55, 72)
    border.line.width = Pt(1.2)


# =====================================
# SMART IMAGE PROMPTS
# =====================================

def build_image_prompt(title, slide_type):
    """
    Prevent repeated photos
    """

    mapping = {
        "problem": "business problem visual",
        "solution": "startup solution illustration",
        "market": "market growth analytics visual",
        "team": "professional startup team",
        "finance": "financial dashboard visual",
        "timeline": "business roadmap visual",
        "comparison": "comparison strategy visual",
        "summary": "executive presentation visual"
    }

    base = mapping.get(
        slide_type,
        "modern business consulting visual"
    )

    return (
        f"{title}, "
        f"{base}, "
        f"premium presentation style, "
        f"minimal professional composition"
    )


# =====================================
# MAIN GENERATOR
# =====================================

def create_pptx(slides_data: list) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    image_service = ImageService()

    # theme from first user title
    main_prompt = slides_data[0].get("title", "")
    theme = resolve_theme_from_prompt(main_prompt)

    bg_color = theme["bg"]
    accent_color = theme["accent"]
    title_color = theme["title"]
    text_color = theme["text"]

    # =====================================
    # ONLY CONTENT SLIDES
    # no cover
    # no final slide
    # =====================================

    for index, slide_data in enumerate(slides_data):
        title = slide_data.get("title", "Без названия")
        content = slide_data.get("content", "")
        slide_type = slide_data.get("slide_type", "bullet")

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # alternating layout
        image_on_left = (index % 2 == 1)

        # if title too long → smaller image
        smaller_image = len(title) > 42

        add_background(
            slide,
            prs,
            bg_color
        )

        add_accent_line(
            slide,
            accent_color
        )

        add_title(
            slide,
            title,
            title_color,
            image_on_left=image_on_left
        )

        add_content(
            slide,
            content,
            text_color,
            image_on_left=image_on_left
        )

        try:
            image_prompt = build_image_prompt(
                title,
                slide_type
            )

            print(f"Searching image: {image_prompt}")

            image_path = image_service.generate_and_download(
                prompt=image_prompt
            )

            add_image_card(
                slide,
                image_path,
                image_on_left=image_on_left,
                smaller=smaller_image
            )
            if os.path.exists(image_path):
                os.remove(image_path)
                print(f"Temporary image deleted: {image_path}")

            print(f"Image added: {image_path}")

        except Exception as e:
            print(f"Image skipped: {str(e)}")

    # =====================================
    # SAVE
    # =====================================

    os.makedirs("generated", exist_ok=True)

    filename = (
        f"generated/presentation_"
        f"{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    )

    prs.save(filename)

    print(f"PPTX saved: {filename}")

    return filename